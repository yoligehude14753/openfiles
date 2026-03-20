from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from .base import BaseParser

class PPTXParser(BaseParser):
    SUPPORTED_TYPES = ["pptx"]

    def can_parse(self, file_type: str) -> bool:
        return file_type in self.SUPPORTED_TYPES

    def parse(self, file_path: Path) -> Dict[str, Any]:
        try:
            prs = Presentation(file_path)

            slides_data = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from shapes
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)

                # Extract notes
                notes_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text

                slide_data = {
                    'page_number': slide_num,
                    'text': '\n'.join(texts),
                    'notes': notes_text,
                    'shape_count': len(slide.shapes)
                }
                slides_data.append(slide_data)

            # Combine all text
            all_text = '\n\n'.join([s['text'] for s in slides_data])

            return {
                'content': all_text,
                'slides': slides_data,
                'num_slides': len(slides_data),
                'success': True
            }
        except Exception as e:
            return {
                'content': '',
                'error': str(e),
                'success': False
            }
